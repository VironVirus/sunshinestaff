from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
LOGO_PATH = ROOT / "public" / "images" / "logo.jpg"
OUTPUT_PATH = (
    ROOT
    / "docs"
    / "Sunshine-Hotel-Front-Office-HouseKeeping-Formal-Projector-Deck.pptx"
)

NAVY = RGBColor(20, 32, 53)
NAVY_SOFT = RGBColor(33, 52, 84)
GOLD = RGBColor(191, 154, 74)
WHITE = RGBColor(255, 255, 255)
IVORY = RGBColor(248, 246, 241)
SLATE = RGBColor(82, 94, 112)
BLACKISH = RGBColor(17, 24, 39)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(220, 38, 38)
BLUE = RGBColor(2, 132, 199)
SOFT_GREEN = RGBColor(236, 253, 245)
SOFT_RED = RGBColor(254, 242, 242)
SOFT_BLUE = RGBColor(240, 249, 255)
SOFT_DARK = RGBColor(241, 245, 249)


def add_logo(slide, left, top, width):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)


def set_runs(paragraph, size, color, bold=False, name="Aptos"):
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = name


def add_footer(slide, slide_number, label="Formal Projector Edition"):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(7.04),
        Inches(12.2),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.08), Inches(8.5), Inches(0.25))
    footer_frame = footer.text_frame
    footer_frame.text = f"Sunshine Hotel Staff Portal | {label} | Powered by Tapxora"
    set_runs(footer_frame.paragraphs[0], 10, SLATE)

    page = slide.shapes.add_textbox(Inches(10.75), Inches(7.05), Inches(1.7), Inches(0.25))
    page_frame = page.text_frame
    page_frame.text = f"Slide {slide_number}"
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_runs(page_frame.paragraphs[0], 10, SLATE)


def add_header(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.333),
        Inches(0.75),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    add_logo(slide, Inches(11.6), Inches(0.08), Inches(1.0))

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.95), Inches(9.5), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    set_runs(title_frame.paragraphs[0], 27, NAVY, bold=True)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.72), Inches(1.42), Inches(10.0), Inches(0.38)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        set_runs(subtitle_frame.paragraphs[0], 12, GOLD, bold=True)


def add_title_slide(prs, slide_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = IVORY
    bg.line.fill.background()

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(0.95),
        Inches(11.9),
        Inches(5.55),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1.1),
        Inches(1.35),
        Inches(0.22),
        Inches(4.65),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    add_logo(slide, Inches(11.1), Inches(1.2), Inches(1.15))

    small = slide.shapes.add_textbox(Inches(1.55), Inches(1.38), Inches(4.5), Inches(0.3))
    small_frame = small.text_frame
    small_frame.text = "FORMAL PROJECTOR EDITION"
    set_runs(small_frame.paragraphs[0], 11, GOLD, bold=True)

    title = slide.shapes.add_textbox(Inches(1.55), Inches(2.0), Inches(8.2), Inches(1.55))
    title_frame = title.text_frame
    title_frame.text = "Sunshine Hotel Staff Portal"
    set_runs(title_frame.paragraphs[0], 31, WHITE, bold=True)
    p = title_frame.add_paragraph()
    p.text = "Training for Front Office and HouseKeeping Leaders"
    set_runs(p, 23, WHITE, bold=True)

    subtitle = slide.shapes.add_textbox(Inches(1.55), Inches(3.85), Inches(8.4), Inches(1.15))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = (
        "Designed for managers and supervisors using live room control, cleaned-room release, "
        "complaints, reports, and shift supervision."
    )
    set_runs(subtitle_frame.paragraphs[0], 18, RGBColor(223, 229, 239))

    audience = slide.shapes.add_textbox(Inches(1.55), Inches(5.2), Inches(8.5), Inches(0.6))
    audience_frame = audience.text_frame
    audience_frame.text = (
        "Audience: Front Office Manager, Front Office Supervisors, "
        "HouseKeeping Manager, HouseKeeping Supervisors"
    )
    set_runs(audience_frame.paragraphs[0], 14, GOLD, bold=True)

    add_footer(slide, slide_number)


def add_section_slide(prs, slide_number, title, subtitle, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    block = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.95),
        Inches(1.15),
        Inches(11.45),
        Inches(4.85),
    )
    block.fill.solid()
    block.fill.fore_color.rgb = NAVY
    block.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.95),
        Inches(1.15),
        Inches(0.28),
        Inches(4.85),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    add_logo(slide, Inches(11.0), Inches(1.35), Inches(1.0))

    tbox = slide.shapes.add_textbox(Inches(1.45), Inches(2.0), Inches(7.6), Inches(1.0))
    tf = tbox.text_frame
    tf.text = title
    set_runs(tf.paragraphs[0], 30, WHITE, bold=True)

    sbox = slide.shapes.add_textbox(Inches(1.45), Inches(3.1), Inches(8.5), Inches(1.25))
    sf = sbox.text_frame
    sf.text = subtitle
    set_runs(sf.paragraphs[0], 18, RGBColor(223, 229, 239))

    add_footer(slide, slide_number)


def add_bullet_slide(prs, slide_number, title, subtitle, bullets, side_title=None, side_points=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, subtitle)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.85),
        Inches(7.8),
        Inches(4.95),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = IVORY
    panel.line.color.rgb = RGBColor(232, 234, 238)

    text_box = slide.shapes.add_textbox(Inches(0.95), Inches(2.1), Inches(7.1), Inches(4.4))
    tf = text_box.text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        set_runs(p, 21, BLACKISH)

    if side_title and side_points:
        side = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.72),
            Inches(1.85),
            Inches(3.85),
            Inches(4.95),
        )
        side.fill.solid()
        side.fill.fore_color.rgb = NAVY_SOFT
        side.line.fill.background()

        st = slide.shapes.add_textbox(Inches(9.0), Inches(2.12), Inches(3.2), Inches(0.4))
        stf = st.text_frame
        stf.text = side_title
        set_runs(stf.paragraphs[0], 16, GOLD, bold=True)

        sb = slide.shapes.add_textbox(Inches(9.0), Inches(2.55), Inches(3.1), Inches(3.8))
        sbf = sb.text_frame
        sbf.clear()
        for idx, point in enumerate(side_points):
            p = sbf.paragraphs[0] if idx == 0 else sbf.add_paragraph()
            p.text = point
            p.level = 0
            p.space_after = Pt(10)
            set_runs(p, 16, WHITE)

    add_footer(slide, slide_number)


def add_two_panel_slide(prs, slide_number, title, subtitle, left_title, left_points, right_title, right_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, subtitle)

    for left, width, heading, points in [
        (Inches(0.68), Inches(5.9), left_title, left_points),
        (Inches(6.75), Inches(5.9), right_title, right_points),
    ]:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            Inches(1.9),
            width,
            Inches(4.88),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = IVORY
        panel.line.color.rgb = RGBColor(232, 234, 238)

        head = slide.shapes.add_textbox(left + Inches(0.25), Inches(2.15), width - Inches(0.5), Inches(0.36))
        head_frame = head.text_frame
        head_frame.text = heading
        set_runs(head_frame.paragraphs[0], 16, NAVY, bold=True)

        body = slide.shapes.add_textbox(left + Inches(0.25), Inches(2.55), width - Inches(0.45), Inches(3.8))
        bf = body.text_frame
        bf.clear()
        for idx, point in enumerate(points):
            p = bf.paragraphs[0] if idx == 0 else bf.add_paragraph()
            p.text = point
            p.level = 0
            p.space_after = Pt(8)
            set_runs(p, 18, BLACKISH)

    add_footer(slide, slide_number)


def add_status_slide(prs, slide_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "HouseKeeping Status Codes", "Status colors that must be understood the same way by every leader")

    statuses = [
        ("Occupied", "Black", "Guest is still staying in the room.", BLACKISH, SOFT_DARK),
        ("Vacant and Cleaned", "Green", "Room is empty and ready for sale.", GREEN, SOFT_GREEN),
        ("Out of Order", "Red", "Room cannot be sold until fixed and cleared.", RED, SOFT_RED),
        ("Vacant and Uncleaned", "Blue", "Room is empty but not ready for sale yet.", BLUE, SOFT_BLUE),
    ]

    positions = [
        (Inches(0.8), Inches(2.0)),
        (Inches(6.9), Inches(2.0)),
        (Inches(0.8), Inches(4.3)),
        (Inches(6.9), Inches(4.3)),
    ]

    for (title, color_name, meaning, border, fill), (left, top) in zip(statuses, positions):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(5.55),
            Inches(1.7),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = border

        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left + Inches(0.22),
            top + Inches(0.22),
            Inches(1.32),
            Inches(0.42),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = border
        badge.line.fill.background()

        badge_text = slide.shapes.add_textbox(left + Inches(0.38), top + Inches(0.28), Inches(1.0), Inches(0.2))
        btf = badge_text.text_frame
        btf.text = color_name
        btf.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_runs(btf.paragraphs[0], 10, WHITE, bold=True)

        title_box = slide.shapes.add_textbox(left + Inches(1.75), top + Inches(0.18), Inches(3.4), Inches(0.3))
        ttf = title_box.text_frame
        ttf.text = title
        set_runs(ttf.paragraphs[0], 16, NAVY, bold=True)

        body_box = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.78), Inches(5.0), Inches(0.7))
        btf2 = body_box.text_frame
        btf2.text = meaning
        set_runs(btf2.paragraphs[0], 15, SLATE)

    add_footer(slide, slide_number)


def add_matrix_slide(prs, slide_number, title, subtitle, columns, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, subtitle)

    table = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        Inches(0.6),
        Inches(1.95),
        Inches(12.1),
        Inches(4.95),
    ).table

    widths = [Inches(2.8), Inches(4.3), Inches(5.0)]
    for idx, width in enumerate(widths[: len(columns)]):
        table.columns[idx].width = width

    for col_idx, heading in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = heading
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            set_runs(paragraph, 11, WHITE, bold=True)

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = IVORY if row_idx % 2 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                set_runs(paragraph, 11, BLACKISH)

    add_footer(slide, slide_number)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Formal Projector Deck - Sunshine Hotel Staff Portal"
    prs.core_properties.subject = "Front Office and HouseKeeping leader onboarding"
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.company = "Tapxora"

    slide_number = 1
    add_title_slide(prs, slide_number)
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Session Purpose",
        "Why this version exists",
        [
            "Train Front Office and HouseKeeping leaders on the exact live workflows they use every shift.",
            "Keep slides projector-friendly with larger text, cleaner branding, and easier spoken delivery.",
            "Focus on speed, accuracy, and cross-department room coordination.",
            "Prepare managers and supervisors to use the app independently after training.",
        ],
        "Use this deck when",
        [
            "Training in a meeting room",
            "Projecting to a group",
            "Running live demonstrations",
            "Leading sign-off drills",
        ],
    )
    slide_number += 1

    add_matrix_slide(
        prs,
        slide_number,
        "Audience and Access",
        "Who this deck is for and what each group can control",
        ["Role", "Main access", "Important limit"],
        [
            ["Front Office Manager", "Full Front Office work tools, reports, complaints, events, team tools.", "Cannot publish cleaned rooms."],
            ["Front Office Supervisor", "Same operational access as Front Office manager for now.", "Cannot publish cleaned rooms."],
            ["HouseKeeping Manager", "Cleaned-room publishing, HouseKeeping reports, complaints, property issues, team tools.", "Does not control Front Office guest check-in/out."],
            ["HouseKeeping Supervisor", "Same operational access as HouseKeeping manager for now.", "Must use room statuses carefully and consistently."],
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Portal Orientation",
        "What leaders should understand before doing any update",
        "Layout",
        [
            "Header shows logo, staff name, department, title, and Sign out.",
            "Summary cards show live operational numbers.",
            "Tabs separate Work, Information, and Events and Bookings.",
            "Mobile uses tab pills for switching sections.",
        ],
        "Mindset",
        [
            "Use Work for operational changes.",
            "Use Information for personal dashboard, announcements, birthdays, and news.",
            "Use Events and Bookings when Front Office needs to create or edit event records.",
            "Do not guess; verify the screen before saving.",
        ],
    )
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Core Operating Rules",
        "These rules protect data quality",
        [
            "Always select Floor before Room.",
            "Publish updates as soon as the room status changes in real life.",
            "At 6:00 AM the hotel operational day resets, but active multi-day bookings continue automatically.",
            "Out-of-order rooms must never be sold until cleared.",
            "Unresolved complaints stay in daily reporting until fixed.",
            "Reports should be printed or downloaded before shift handover.",
        ],
        "Trainer phrase",
        [
            "If the room changed in real life, it must change in the app immediately."
        ],
    )
    slide_number += 1

    add_section_slide(
        prs,
        slide_number,
        "Front Office Leadership",
        "Check-in, check-out, room move, complaints, and reporting",
        GOLD,
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Front Office Shift Opening Routine",
        "Front office",
        "Check at the start of shift",
        [
            "Confirm your account name, department, and title are correct.",
            "Review In-house, Available rooms, Breakfast entitlement, and Cleaned rooms.",
            "Open the HouseKeeping Cleaned Rooms Report to see ready rooms.",
            "Review unresolved complaints before allocating rooms.",
        ],
        "Why this matters",
        [
            "Prevents selling bad rooms.",
            "Keeps breakfast numbers accurate.",
            "Reduces disagreement between Front Office and HouseKeeping.",
            "Improves handover quality.",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Front Office Workflow: Check In",
        "Front office",
        "Steps",
        [
            "Open Work > Rooms > Check in guest.",
            "Select Floor and Room.",
            "Set Booked for days.",
            "Set Breakfast Yes/No and Breakfast count if needed.",
            "Click Check in room and wait for success confirmation.",
        ],
        "What should update",
        [
            "In-house increases.",
            "Available rooms reduces.",
            "Breakfast entitlement updates.",
            "Room appears in the In House Report.",
            "Freshly cleaned tag disappears if that room was just released by HouseKeeping.",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Front Office Workflow: Check Out and Room Move",
        "Front office",
        "Check out",
        [
            "Open Check out room form.",
            "Select floor and occupied room.",
            "Click Mark checked out.",
            "Confirm the room leaves occupied stock.",
        ],
        "Room move",
        [
            "Use Room move when the guest changes rooms but remains in house.",
            "Select current occupied room, then select the new unoccupied room.",
            "Click Move guest.",
            "The move is saved into the daily report for that operational date.",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Front Office Reports",
        "Front office",
        "Available reports",
        [
            "In House Report",
            "HouseKeeping Cleaned Rooms Report",
            "Daily Report",
            "Date Range Report",
        ],
        "What daily/date-range reporting now shows",
        [
            "Room counts and breakfast entitlement",
            "Events for the selected day",
            "Room moves for the selected day",
            "Complaint follow-up until resolution",
            "Fixed status on the resolution day only",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Front Office Complaints, Events, and Team Tools",
        "Front office",
        "Complaints",
        [
            "Log issues with exact room, complaint type, and clear note.",
            "Clear a complaint only when it is actually resolved.",
            "Remember: unresolved complaints continue into later daily reports.",
        ],
        "Other tools",
        [
            "Events and Bookings: create and update event records.",
            "Team: assign and remove shift dates for staff.",
            "Information: news, announcements, birthdays, recognition, and personal details.",
        ],
    )
    slide_number += 1

    add_section_slide(
        prs,
        slide_number,
        "HouseKeeping Leadership",
        "Cleaned-room release, room-status reports, and property escalation",
        GREEN,
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "HouseKeeping Shift Opening Routine",
        "Housekeeping",
        "Check at the start of shift",
        [
            "Review In-house, Available rooms, and Cleaned rooms summary cards.",
            "Open the cleaned-room board and see what was already released.",
            "Check Morning or Afternoon report status before new entries.",
            "Review active complaints and out-of-order rooms.",
        ],
        "Main aim",
        [
            "Keep room readiness honest and current.",
            "Release clean rooms quickly.",
            "Prevent Front Office from selling rooms that are not truly ready.",
            "Support handover with visible reports instead of memory.",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "HouseKeeping Workflow: Publish Cleaned Rooms",
        "Housekeeping",
        "How to do it",
        [
            "Open Work > Rooms > Mark cleaned room.",
            "Select Floor and Room.",
            "Click Publish cleaned room.",
            "Check the cleaned-room board to confirm the room appears.",
        ],
        "Important limit",
        [
            "Only HouseKeeping leaders publish cleaned rooms.",
            "Front Office leaders can view and report on cleaned rooms but cannot publish them.",
            "Clear any wrong cleaned-room entry immediately.",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "HouseKeeping Workflow: Morning and Afternoon Reports",
        "Housekeeping",
        "How to complete a report",
        [
            "Open Work > HouseKeeping Reports.",
            "Choose Morning report or Afternoon report.",
            "Select Floor, Room, and Status.",
            "Save each room and review the board on the right.",
            "Print or download once the round is complete.",
        ],
        "What supervisors should watch",
        [
            "Use the correct status every time.",
            "Keep the board updated room by room.",
            "Review selected room details before clearing an entry.",
            "Use the report as the internal truth for that shift.",
        ],
    )
    slide_number += 1

    add_status_slide(prs, slide_number)
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "HouseKeeping Property, Complaints, and Team Tools",
        "Housekeeping",
        "Property and complaints",
        [
            "Use Property to mark rooms Out of order and explain what must be done.",
            "Use Complaints to review and clear issues after real resolution.",
            "Do not clear a room issue early just to make stock look better.",
        ],
        "Team and communication tools",
        [
            "Assign shifts in the Team section.",
            "Remove wrong shifts quickly.",
            "Use Information for announcements, news, birthdays, and staff details.",
            "Review Events and Bookings when occupancy pressure may affect cleaning priorities.",
        ],
    )
    slide_number += 1

    add_section_slide(
        prs,
        slide_number,
        "Joint Coordination",
        "How Front Office and HouseKeeping should work together inside the app",
        BLUE,
    )
    slide_number += 1

    add_matrix_slide(
        prs,
        slide_number,
        "Cross-Department Workflow",
        "The intended room life cycle inside the portal",
        ["Stage", "Who acts", "Expected result"],
        [
            ["Guest checks in", "Front Office", "Room becomes occupied and appears in in-house reporting."],
            ["Guest checks out", "Front Office", "Room leaves occupied stock and becomes vacant."],
            ["Room cleaned", "HouseKeeping", "Room appears on cleaned-room board and as freshly cleaned for Front Office."],
            ["Room sold again", "Front Office", "Freshly cleaned tag disappears because room is occupied."],
            ["Room has fault", "HouseKeeping or Maintainance", "Room is marked out of order and removed from sellable stock."],
            ["Fault fixed", "HouseKeeping or Maintainance", "Room is cleared and can return to normal availability."],
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Practical Assessment Drills",
        "Trainer-led exercises",
        "Front Office drill",
        [
            "Check in one room with breakfast.",
            "Check out one room correctly.",
            "Move one guest to another room.",
            "Print the daily report.",
            "Run and download a date-range report.",
        ],
        "HouseKeeping drill",
        [
            "Publish three cleaned rooms.",
            "Complete one morning or afternoon report with mixed statuses.",
            "Mark one room out of order with proper note.",
            "Print the HouseKeeping report.",
            "Clear one wrong entry correctly.",
        ],
    )
    slide_number += 1

    add_two_panel_slide(
        prs,
        slide_number,
        "Common Mistakes to Correct Early",
        "Trainer watchpoints",
        "Front Office errors",
        [
            "Wrong floor selected before room choice",
            "Breakfast count entered wrongly",
            "Guest not checked out on time",
            "Room move done without verifying destination readiness",
        ],
        "HouseKeeping errors",
        [
            "Publishing a room before it is truly clean",
            "Using the wrong status color meaning",
            "Clearing out-of-order too early",
            "Leaving reports incomplete before handover",
        ],
    )
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Go-Live Sign-Off",
        "Minimum readiness standard",
        [
            "Every leader must log in without help.",
            "Front Office leaders must complete one check-in, one check-out, one room move, and one report download alone.",
            "HouseKeeping leaders must complete one cleaned-room publish, one report update, and one out-of-order update alone.",
            "Both departments must explain the rule for complaint carry-over and report follow-up.",
            "Both departments must agree on who can publish cleaned rooms and who can only view them.",
        ],
        "Escalation route",
        [
            "1. Verify account",
            "2. Refresh and retry",
            "3. Capture room/date details",
            "4. Escalate to IT admin",
        ],
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    result = build_presentation()
    print(result)
