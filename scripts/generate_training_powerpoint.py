from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
LOGO_PATH = ROOT / "public" / "images" / "logo.jpg"
OUTPUT_PATH = ROOT / "docs" / "Sunshine-Hotel-Front-Office-Housekeeping-Onboarding-Training.pptx"

NAVY = RGBColor(22, 35, 56)
GOLD = RGBColor(197, 157, 64)
SLATE = RGBColor(71, 85, 105)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(5, 150, 105)
RED = RGBColor(220, 38, 38)
BLUE = RGBColor(2, 132, 199)
BLACKISH = RGBColor(15, 23, 42)


def add_logo(slide, left, top, width):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)


def style_text_frame(text_frame, font_size=24, color=BLACKISH, bold=False):
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    for paragraph in text_frame.paragraphs:
        paragraph.space_after = Pt(8)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Aptos"


def add_footer(slide, slide_number):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.55),
        Inches(7.05),
        Inches(12.2),
        Inches(0.03),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.65), Inches(7.08), Inches(9.5), Inches(0.25))
    footer_frame = footer.text_frame
    footer_frame.text = "Sunshine Hotel Staff Portal Training | Powered by Tapxora"
    style_text_frame(footer_frame, font_size=10, color=SLATE)

    page_box = slide.shapes.add_textbox(Inches(10.8), Inches(7.04), Inches(1.6), Inches(0.3))
    page_frame = page_box.text_frame
    page_frame.text = f"Slide {slide_number}"
    style_text_frame(page_frame, font_size=10, color=SLATE)
    page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_header(slide, title, eyebrow=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(13.333),
        Inches(0.72),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    if eyebrow:
        eyebrow_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.2), Inches(3.0), Inches(0.22))
        eyebrow_frame = eyebrow_box.text_frame
        eyebrow_frame.text = eyebrow.upper()
        style_text_frame(eyebrow_frame, font_size=10, color=GOLD, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.88), Inches(8.7), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    style_text_frame(title_frame, font_size=26, color=NAVY, bold=True)

    add_logo(slide, Inches(11.55), Inches(0.08), Inches(1.05))


def add_callout(slide, title, body, left, top, width, height, fill_color=LIGHT):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = WHITE

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.15), width - Inches(0.35), Inches(0.35))
    title_frame = title_box.text_frame
    title_frame.text = title
    style_text_frame(title_frame, font_size=14, color=NAVY, bold=True)

    body_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.52), width - Inches(0.35), height - Inches(0.68))
    body_frame = body_box.text_frame
    body_frame.text = body
    style_text_frame(body_frame, font_size=12, color=SLATE)


def add_bullets(text_frame, items, font_size=20, color=BLACKISH):
    text_frame.clear()
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        paragraph.text = text
        paragraph.level = level
        paragraph.space_after = Pt(8 if level == 0 else 4)
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_size if level == 0 else font_size - 2)
            run.font.color.rgb = color
            run.font.name = "Aptos"


def add_title_slide(prs, slide_number):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(0.75),
        Inches(12.1),
        Inches(5.7),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = NAVY
    hero.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.95),
        Inches(1.15),
        Inches(0.18),
        Inches(4.9),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()

    add_logo(slide, Inches(10.95), Inches(1.05), Inches(1.2))

    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(1.3), Inches(8.6), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Sunshine Hotel Staff Portal"
    style_text_frame(title_frame, font_size=30, color=WHITE, bold=True)
    p = title_frame.add_paragraph()
    p.text = "Front Office and HouseKeeping Manager/Supervisor Onboarding"
    p.space_before = Pt(8)
    for run in p.runs:
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = "Aptos"

    subtitle_box = slide.shapes.add_textbox(Inches(1.4), Inches(3.0), Inches(7.5), Inches(1.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = (
        "Detailed practical training for live room control, report handling, complaints, "
        "cleaned-room coordination, and shift supervision."
    )
    style_text_frame(subtitle_frame, font_size=18, color=RGBColor(220, 226, 235))

    detail_box = slide.shapes.add_textbox(Inches(1.4), Inches(4.95), Inches(8.0), Inches(0.9))
    detail_frame = detail_box.text_frame
    detail_frame.text = "Audience: Front Office Manager, Front Office Supervisors, HouseKeeping Manager, HouseKeeping Supervisors"
    style_text_frame(detail_frame, font_size=14, color=GOLD, bold=True)

    add_footer(slide, slide_number)


def add_bullet_slide(prs, slide_number, title, eyebrow, bullets, sidebar_title=None, sidebar_body=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, eyebrow)

    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(8.0), Inches(5.15))
    body_frame = body.text_frame
    add_bullets(body_frame, bullets)

    if sidebar_title and sidebar_body:
        add_callout(
            slide,
            sidebar_title,
            sidebar_body,
            Inches(9.05),
            Inches(1.6),
            Inches(3.45),
            Inches(4.6),
            fill_color=RGBColor(244, 238, 224),
        )

    add_footer(slide, slide_number)


def add_two_column_slide(prs, slide_number, title, eyebrow, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, eyebrow)

    left_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.55),
        Inches(5.95),
        Inches(5.2),
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = LIGHT
    left_panel.line.fill.background()

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.78),
        Inches(1.55),
        Inches(5.9),
        Inches(5.2),
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = LIGHT
    right_panel.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.78), Inches(5.2), Inches(0.35))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    style_text_frame(left_title_frame, font_size=16, color=NAVY, bold=True)

    right_title_box = slide.shapes.add_textbox(Inches(7.02), Inches(1.78), Inches(5.1), Inches(0.35))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    style_text_frame(right_title_frame, font_size=16, color=NAVY, bold=True)

    left_body = slide.shapes.add_textbox(Inches(0.9), Inches(2.15), Inches(5.15), Inches(4.2))
    add_bullets(left_body.text_frame, left_items, font_size=18)

    right_body = slide.shapes.add_textbox(Inches(7.02), Inches(2.15), Inches(5.05), Inches(4.2))
    add_bullets(right_body.text_frame, right_items, font_size=18)

    add_footer(slide, slide_number)


def add_matrix_slide(prs, slide_number, title, eyebrow, columns, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, title, eyebrow)

    table = slide.shapes.add_table(
        len(rows) + 1,
        len(columns),
        Inches(0.7),
        Inches(1.6),
        Inches(11.95),
        Inches(5.4),
    ).table

    table.columns[0].width = Inches(3.0)
    for index in range(1, len(columns)):
        table.columns[index].width = Inches(4.45)

    for col_idx, heading in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = heading
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.size = Pt(12)
                run.font.name = "Aptos"

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if row_idx % 2 == 1 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = BLACKISH
                    run.font.name = "Aptos"

    add_footer(slide, slide_number)


def add_section_slide(prs, slide_number, eyebrow, title, subtitle, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.0),
        Inches(11.8),
        Inches(4.9),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.8),
        Inches(1.0),
        Inches(0.3),
        Inches(4.9),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    add_logo(slide, Inches(10.95), Inches(1.3), Inches(1.05))

    eyebrow_box = slide.shapes.add_textbox(Inches(1.35), Inches(1.45), Inches(4.0), Inches(0.3))
    eyebrow_frame = eyebrow_box.text_frame
    eyebrow_frame.text = eyebrow.upper()
    style_text_frame(eyebrow_frame, font_size=11, color=accent_color, bold=True)

    title_box = slide.shapes.add_textbox(Inches(1.35), Inches(2.0), Inches(8.0), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = title
    style_text_frame(title_frame, font_size=28, color=WHITE, bold=True)

    subtitle_box = slide.shapes.add_textbox(Inches(1.35), Inches(3.2), Inches(8.6), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    style_text_frame(subtitle_frame, font_size=18, color=RGBColor(220, 226, 235))

    add_footer(slide, slide_number)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Sunshine Hotel Staff Portal Onboarding Training"
    prs.core_properties.subject = "Front Office and HouseKeeping manager/supervisor training"
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.company = "Tapxora"

    slide_number = 1
    add_title_slide(prs, slide_number)
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Training Agenda",
        "Session map",
        [
            "1. Understand exactly who can do what in the portal.",
            "2. Learn the layout of the Work, Information, and Events and Bookings tabs.",
            "3. Train Front Office leaders on check-in, check-out, room move, complaints, and reports.",
            "4. Train HouseKeeping leaders on cleaned rooms, report boards, out-of-order rooms, and coordination.",
            "5. Finish with practice drills, data-quality rules, and go-live sign-off.",
        ],
        "Recommended Format",
        "Use this deck in a guided session. Demonstrate each task live in the app, then let every participant repeat the task alone before moving on.",
    )
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Learning Outcomes",
        "Expected result",
        [
            "Log in successfully and identify the correct dashboard for each role.",
            "Use the correct tab and work section without guessing.",
            "Front Office leaders can check in, check out, and move guests correctly.",
            "HouseKeeping leaders can publish cleaned rooms and complete morning/afternoon room reports correctly.",
            "Both departments can handle complaints, reports, and shift supervision with confidence.",
            "Both departments understand which updates are real-time and why speed matters.",
        ],
        "Important Reminder",
        "Managers and supervisors currently use the same operational tools. Line staff do not use these work boards; they only use their staff dashboard and information area.",
    )
    slide_number += 1

    add_matrix_slide(
        prs,
        slide_number,
        "Access and Capability Matrix",
        "Role control",
        [
            "Capability",
            "Front Office Manager/Supervisor",
            "HouseKeeping Manager/Supervisor",
        ],
        [
            ["Work > Rooms", "Can check in, check out, move guests, and print/download reports.", "Can view room metrics and cleaned-room board."],
            ["Publish cleaned rooms", "View only. Front Office must not publish cleaned rooms.", "Yes. Publish and clear cleaned-room entries."],
            ["HouseKeeping Reports", "View not primary in current workflow.", "Yes. Morning and afternoon reports, print, and download."],
            ["Complaints", "Can report and clear room complaints.", "Can report and clear room complaints."],
            ["Property issues", "View only.", "Can mark rooms out of order and clear them after fix."],
            ["Events and Bookings", "Can create and update events/bookings.", "Can view events/bookings for coordination."],
            ["Team", "Can assign and remove shift dates for team members.", "Can assign and remove shift dates for team members."],
            ["Information", "View news, announcements, birthdays, recognition, and personal details.", "View news, announcements, birthdays, recognition, and personal details."],
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Portal Layout and Navigation",
        "Screen map",
        "What every participant should notice",
        [
            "Header shows logo, staff name, department, title, and Sign out.",
            "Summary cards show live numbers such as In-house, Available rooms, Cleaned rooms, and Team.",
            "Top tabs separate Work, Information, and Events and Bookings.",
            "Inside Work, each role sees only the sections they are allowed to use.",
        ],
        "How to move through the app",
        [
            "Start in Work for operational tasks.",
            "Use Information for personal dashboard, shifts, news, announcements, birthdays, and recognition.",
            "Front Office leaders use Events and Bookings for event entries.",
            "On mobile, use the tab pills at the top to switch sections instead of scrolling aimlessly.",
        ],
    )
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Non-Negotiable Operating Rules",
        "Data discipline",
        [
            "Always choose Floor first, then Room. Never guess a room number and never rely on memory alone.",
            "Publish changes as soon as the real room status changes. Delayed entry creates sales, cleaning, and breakfast errors.",
            "At 6:00 AM hotel time the operational day rolls over. Multi-day bookings stay active because remaining days reduce automatically.",
            "Out-of-order rooms must not be assigned by Front Office until HouseKeeping/Maintainance clears them.",
            "Unresolved room complaints continue to appear in daily reports until they are fixed.",
            "Generate reports before handover so the next shift sees the same story you saw.",
        ],
        "Trainer language to use",
        "Repeat this often: 'If the room changed in real life, it must change in the app immediately.'",
    )
    slide_number += 1

    add_section_slide(
        prs,
        slide_number,
        "Front Office Module",
        "Front Office Manager and Supervisor Training",
        "This section trains live room control, guest movement, reporting, and complaint follow-up for Front Office leaders.",
        GOLD,
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Front Office Role Overview",
        "Front office",
        "Main responsibilities",
        [
            "Control room sales status inside Work > Rooms.",
            "Maintain accurate in-house rooms, available rooms, and breakfast entitlement.",
            "Move guests between rooms correctly when room changes happen.",
            "Raise complaints quickly and keep follow-up visible in reports.",
            "Prepare printable daily and date-range reports for review and handover.",
        ],
        "Start-of-shift checklist",
        [
            "Log in and confirm your name, department, and title are correct.",
            "Check the summary cards for In-house, Available rooms, Breakfast entitlement, and Cleaned rooms.",
            "Open the HouseKeeping Cleaned Rooms Report and review rooms tagged freshly cleaned.",
            "Check active complaints before beginning new room sales.",
            "Check Events and Bookings if there is any event affecting rooms, breakfast, or guest movement.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Front Office Workflow: Check In Guest",
        "Front office",
        "How to do it",
        [
            "Open Work > Rooms and locate the Check in guest form.",
            "Select Floor first.",
            "Select Room from the dropdown. Rooms marked '- freshly cleaned' should be preferred because HouseKeeping has already released them.",
            "Choose Booked for and enter the number of days.",
            "Choose Breakfast Yes or No. If Yes, enter Breakfast count.",
            "Click Check in room and wait for the success message before moving to another task.",
        ],
        "What should change immediately",
        [
            "In-house increases.",
            "Available rooms reduces.",
            "Breakfast entitlement updates.",
            "That room appears in the printable In House Report.",
            "If the room was freshly cleaned, the cleaned tag disappears because the room is now occupied.",
            "If the room does not appear, check whether it was already occupied or out of order.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Front Office Workflow: Check Out and Room Move",
        "Front office",
        "Check out room",
        [
            "Open the Check out room form in Work > Rooms.",
            "Select the Floor that currently contains the occupied room.",
            "Select the Occupied room from the dropdown.",
            "Click Mark checked out and wait for the success message.",
            "Confirm the room disappears from occupied stock and becomes available again unless another issue blocks it.",
        ],
        "Room move",
        [
            "Use Room move when the guest changes room but remains in house.",
            "Select Current floor and Occupied room.",
            "Select New floor and Unoccupied room.",
            "Click Move guest.",
            "The move is recorded in the daily report for that operational day.",
            "If the destination room shows '- freshly cleaned', that is a good sign the room is ready for use.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Front Office Reports",
        "Front office",
        "Report types in Rooms",
        [
            "In House Report: occupied rooms floor by floor plus breakfast totals.",
            "HouseKeeping Cleaned Rooms Report: cleaned rooms published by HouseKeeping. Front Office can print/download it but cannot publish cleaned rooms.",
            "Daily Report: current operational-day report for rooms, events, room moves, and complaint follow-up.",
            "Date Range Report: manager/supervisor selects From and To dates, then prints or downloads the combined report.",
        ],
        "What the daily report now includes",
        [
            "In-house rooms, available rooms, breakfast entitlement, and cleaned rooms for the selected day.",
            "Events for that day.",
            "Room moves for that day.",
            "Complaint follow-up across days until resolution.",
            "If a complaint is fixed on the selected day, it appears as Fixed on that day and disappears from the next day onward.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Front Office Complaints, Events, and Team Tools",
        "Front office",
        "Complaints workflow",
        [
            "Open Work > Complaints.",
            "Select Floor, Room, Complaint type, and Note, then send the complaint.",
            "Use clear only when the issue has truly been resolved.",
            "Remember: unresolved complaints keep appearing in daily reports until fix day.",
            "Use precise notes such as 'AC not cooling' instead of vague notes such as 'bad room'.",
        ],
        "Other Front Office manager/supervisor tools",
        [
            "Events and Bookings tab: create and update event details for other managers to view.",
            "Team section: assign and remove shift dates for team members.",
            "Information tab: review announcements, birthdays, hotel news, and personal staff details.",
            "Use these tools at handover so the next shift inherits the correct picture.",
        ],
    )
    slide_number += 1

    add_section_slide(
        prs,
        slide_number,
        "HouseKeeping Module",
        "HouseKeeping Manager and Supervisor Training",
        "This section trains cleaned-room publishing, room-status reporting, out-of-order updates, and coordination with Front Office.",
        GREEN,
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "HouseKeeping Role Overview",
        "Housekeeping",
        "Main responsibilities",
        [
            "Release cleaned rooms quickly so Front Office can sell the correct stock.",
            "Maintain room-status accuracy in HouseKeeping Reports.",
            "Mark rooms out of order when they are not saleable and add clear issue notes.",
            "Track complaints and coordinate with Front Office and Maintainance.",
            "Assign and manage team shift dates.",
        ],
        "Start-of-shift checklist",
        [
            "Log in and confirm profile details are correct.",
            "Review In-house, Available rooms, and Cleaned rooms summary cards.",
            "Open Work > Rooms and check the cleaned-room board from the last shift.",
            "Open HouseKeeping Reports and confirm whether you are working Morning report or Afternoon report.",
            "Review active complaints and out-of-order rooms before assigning physical work.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "HouseKeeping Workflow: Publish Cleaned Rooms",
        "Housekeeping",
        "How to publish cleaned rooms",
        [
            "Open Work > Rooms.",
            "Use the Mark cleaned room form.",
            "Select Floor first, then Room.",
            "Click Publish cleaned room and wait for the success message.",
            "The cleaned room appears on the Freshly cleaned rooms board, grouped by floor.",
            "Front Office now sees that room tagged '- freshly cleaned' in their room dropdown.",
        ],
        "Important control rules",
        [
            "Only HouseKeeping manager/supervisor should publish cleaned rooms.",
            "Do not publish a room until it is physically ready for sale.",
            "If a room was published in error, clear it immediately from the cleaned-room board.",
            "Use the board to review which cleaned rooms are already waiting for Front Office.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "HouseKeeping Workflow: Morning and Afternoon Reports",
        "Housekeeping",
        "How to use HouseKeeping Reports",
        [
            "Open Work > HouseKeeping Reports.",
            "Choose Morning report or Afternoon report.",
            "Select Floor, then Room.",
            "Choose the correct Status and click Save morning report or Save afternoon report.",
            "Use the compact room board on the right to review reported rooms floor by floor.",
            "Use Print report or Download report when management needs the report physically or as PDF.",
        ],
        "Why this report matters",
        [
            "It is the internal operational truth for HouseKeeping checks.",
            "It shows room-by-room condition for morning and afternoon rounds.",
            "It prevents long handwritten sheets and makes the status visible immediately.",
            "It helps HouseKeeping supervisors verify rooms without re-walking every floor blindly.",
        ],
    )
    slide_number += 1

    add_matrix_slide(
        prs,
        slide_number,
        "HouseKeeping Status Meanings and Property Rules",
        "Housekeeping",
        ["Status / Rule", "Meaning", "How to use it correctly"],
        [
            ["Occupied - black", "Guest is currently in the room.", "Use when the room is sold and the guest is still staying there."],
            ["Vacant and Cleaned - green", "Room is empty and ready for sale.", "Use only after physical cleaning is completed and checked."],
            ["Out of Order - red", "Room cannot be sold right now.", "Use when damage or a serious issue blocks sale. Add issue note in Property."],
            ["Vacant and Uncleaned - blue", "Room is empty but not ready for sale.", "Use when the room is not occupied but still needs cleaning or reset."],
            ["Property > Out of order", "Room is removed from Front Office sellable stock.", "Select Floor and Room, keep Out of order = Yes, write exactly what needs to be done, then save."],
            ["Clearing a property issue", "Room becomes sellable again after real fix.", "Clear only when the room is truly ready; otherwise Front Office may sell a bad room."],
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "HouseKeeping Complaints, Coordination, and Team Tools",
        "Housekeeping",
        "Complaints and coordination",
        [
            "HouseKeeping leaders can also use Work > Complaints to review and clear room complaints.",
            "When Front Office reports an issue, HouseKeeping should confirm whether it is cleaning-related, room-condition-related, or requires Maintainance.",
            "Use exact notes and communicate resolution as soon as the room is ready.",
            "Remember: until the issue is cleared, it continues to appear in daily report follow-up.",
        ],
        "Team and communication tools",
        [
            "Use Work > Team to assign shift dates for attendants/supervisors.",
            "Remove wrong shift entries immediately to avoid staff confusion.",
            "Use Information for announcements, birthdays, hotel news, and personal staff details.",
            "Use Events and Bookings view to prepare for occupancy pressure, VIP movement, or event-related cleaning needs.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Practical Drills and Assessment",
        "Trainer exercises",
        "Front Office drills",
        [
            "Check in room 104 for 3 days with breakfast count 2.",
            "Check out one occupied room correctly.",
            "Move a guest from one occupied room to a clean available room.",
            "Print the daily report and explain where room moves and complaints appear.",
            "Run a date-range report for two operational dates and download it as PDF.",
        ],
        "HouseKeeping drills",
        [
            "Publish three cleaned rooms from different floors.",
            "Create one Morning report with at least four rooms and mixed statuses.",
            "Mark one room out of order and write the issue note clearly.",
            "Clear one wrong cleaned-room entry.",
            "Print the HouseKeeping report and explain the meaning of each status color.",
        ],
    )
    slide_number += 1

    add_two_column_slide(
        prs,
        slide_number,
        "Common Mistakes and How to Correct Them",
        "Risk control",
        "Front Office mistakes",
        [
            "Selecting the wrong floor before selecting room. Correction: always confirm floor first.",
            "Entering breakfast count when breakfast is No. Correction: leave breakfast count disabled unless breakfast is Yes.",
            "Forgetting to check out a departed guest. Correction: close rooms immediately at departure.",
            "Moving a guest without checking destination room readiness. Correction: prefer rooms tagged freshly cleaned or otherwise confirmed ready.",
        ],
        "HouseKeeping mistakes",
        [
            "Publishing a cleaned room before physical completion. Correction: publish only after inspection.",
            "Using wrong status color meaning. Correction: black occupied, green vacant and cleaned, red out of order, blue vacant and uncleaned.",
            "Clearing out-of-order too early. Correction: clear only after the room is truly ready for sale.",
            "Leaving morning/afternoon report incomplete. Correction: report every relevant room before printing.",
        ],
    )
    slide_number += 1

    add_bullet_slide(
        prs,
        slide_number,
        "Go-Live Checklist and Support",
        "Readiness",
        [
            "Confirm every manager/supervisor can log in alone and reach the correct work section.",
            "Confirm Front Office leaders can complete one check-in, one check-out, one room move, and one daily report download without help.",
            "Confirm HouseKeeping leaders can publish cleaned rooms, complete one HouseKeeping report, and mark one room out of order without help.",
            "Confirm both departments understand who publishes cleaned rooms and who only views them.",
            "Confirm the printer/PDF workflow works on at least one laptop and one phone.",
            "Escalate access problems, missing tabs, sync delays, or broken report downloads to the IT admin immediately.",
        ],
        "Suggested escalation order",
        "1. Verify the user signed into the correct account. 2. Refresh and retry the exact task. 3. Capture a screenshot plus the exact room number/report date. 4. Escalate to IT with those details.",
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    result = build_presentation()
    print(result)
